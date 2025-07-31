import os
import fitz  # PyMuPDF
from docx import Document
from celery import shared_task
from PIL import Image
from django.conf import settings
import pandas as pd
from pptx import Presentation
from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import letter
from pdf2image import convert_from_path

COMPRESSED_DIR = os.path.join(settings.MEDIA_ROOT, 'compressed')
CONVERTED_DIR = os.path.join(settings.MEDIA_ROOT, 'converted')
os.makedirs(COMPRESSED_DIR, exist_ok=True)
os.makedirs(CONVERTED_DIR, exist_ok=True)

@shared_task
def compress_file_task(input_path, output_filename, quality=85, max_width=1920):
    """
    Compresses a file (image, PDF, DOCX) and saves it.

    :param input_path: The absolute path to the uploaded file.
    :param output_filename: The desired name for the output file.
    :param quality: The quality for the output image (1-95).
    :param max_width: The maximum width for the image.
    :return: The path to the compressed file relative to the MEDIA_ROOT.
    """
    try:
        file_extension = os.path.splitext(input_path)[1].lower()
        output_path_full = os.path.join(COMPRESSED_DIR, output_filename)

        if file_extension in ['.jpg', '.jpeg', '.png', '.webp']:
            with Image.open(input_path) as img:
                if img.width > max_width:
                    w_percent = (max_width / float(img.width))
                    h_size = int((float(img.height) * float(w_percent)))
                    img = img.resize((max_width, h_size), Image.Resampling.LANCZOS)
                if img.mode == 'RGBA':
                    img = img.convert('RGB')
                img.save(output_path_full, 'JPEG', optimize=True, quality=quality)
        elif file_extension == '.pdf':
            doc = fitz.open(input_path)
            doc.save(output_path_full, garbage=4, deflate=True, clean=True)
        elif file_extension == '.docx':
            # DOCX are zip files, so we can't do much more than just copying it.
            # For a real-world scenario, we might want to compress images inside the docx.
            import shutil
            shutil.copy(input_path, output_path_full)
        else:
            raise ValueError(f"Unsupported file type for compression: {file_extension}")

        os.remove(input_path)
        return os.path.join('compressed', output_filename)

    except Exception as e:
        if os.path.exists(input_path):
            os.remove(input_path)
        raise e

@shared_task
def convert_file_task(input_path, output_filename, target_format):
    """
    Converts a file to a different format.

    :param input_path: The absolute path to the uploaded file.
    :param output_filename: The desired name for the output file.
    :param target_format: The target format (e.g., 'PDF', 'DOCX', 'PNG').
    :return: The path to the converted file relative to the MEDIA_ROOT.
    """
    try:
        output_path_full = os.path.join(CONVERTED_DIR, output_filename)
        file_extension = os.path.splitext(input_path)[1].lower()

        # Image conversions
        if file_extension in ['.jpg', '.jpeg', '.png', '.webp']:
            with Image.open(input_path) as img:
                if target_format == 'PDF':
                    img.save(output_path_full, 'PDF', resolution=100.0)
                else:
                    img.save(output_path_full, target_format)

        # PDF conversions
        elif file_extension == '.pdf':
            if target_format in ['PNG', 'JPEG']:
                pages = convert_from_path(input_path, 200)
                if pages:
                    pages[0].save(output_path_full, target_format)
            elif target_format == 'DOCX':
                pdf_document = fitz.open(input_path)
                word_document = Document()
                for page_num in range(len(pdf_document)):
                    page = pdf_document.load_page(page_num)
                    text = page.get_text("text")
                    word_document.add_paragraph(text)
                    if page_num < len(pdf_document) - 1:
                        word_document.add_page_break()
                word_document.save(output_path_full)

        # DOCX conversions
        elif file_extension == '.docx':
            if target_format == 'PDF':
                import pypandoc
                pypandoc.convert_file(input_path, 'pdf', outputfile=output_path_full)

        # Spreadsheet conversions
        elif file_extension in ['.xlsx', '.xls']:
            df = pd.read_excel(input_path)
            if target_format == 'CSV':
                df.to_csv(output_path_full, index=False)
            elif target_format == 'XLSX':
                df.to_excel(output_path_full, index=False)

        # Presentation conversions
        elif file_extension in ['.pptx', '.ppt']:
            if target_format == 'PDF':
                prs = Presentation(input_path)
                c = canvas.Canvas(output_path_full, pagesize=letter)
                width, height = letter
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            c.drawString(100, height - 100, shape.text)
                    c.showPage()
                c.save()

        else:
            raise ValueError(f"Unsupported conversion from {file_extension} to {target_format}")

        os.remove(input_path)
        return os.path.join('converted', output_filename)

    except Exception as e:
        if os.path.exists(input_path):
            os.remove(input_path)
        raise e